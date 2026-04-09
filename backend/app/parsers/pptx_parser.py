"""
PowerPoint parser with optional python-pptx dependency.
"""

from io import BytesIO
from typing import List

from .base import BaseParser, DocumentMetadata, DocumentType, ParseResult
from .image_parser import ImageParser, ParseImage
from logger import logger


class PptxParser(BaseParser):
    """Parse PPTX slides and OCR embedded images."""

    def __init__(self) -> None:
        self.image_parser = ImageParser()

    @property
    def supported_extensions(self) -> list:
        return [".pptx"]

    @property
    def parser_name(self) -> str:
        return "PptxParser"

    def can_parse(self, filename: str, mime_type: str = None) -> bool:
        normalized_name = str(filename or "").lower()
        return normalized_name.endswith(".pptx") or (
            mime_type and "presentation" in mime_type.lower()
        )

    def parse(self, file_data: bytes, filename: str) -> ParseResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            logger.error("PPTX parser dependency missing for %s: %s", filename, exc)
            return self._empty_result(len(file_data))

        try:
            presentation = Presentation(BytesIO(file_data))

            parts: List[str] = []
            images: List[ParseImage] = []

            for slide_index, slide in enumerate(presentation.slides):
                slide_texts: List[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)
                    if shape.shape_type == 13:
                        try:
                            image_data = shape.image.blob
                            images.append(
                                ParseImage(
                                    id=self.image_parser.get_image_hash(image_data),
                                    index=len(images),
                                    data=image_data,
                                )
                            )
                        except Exception:
                            continue

                if slide_texts:
                    parts.append(f"--- Slide {slide_index + 1} ---\n" + "\n".join(slide_texts))

            if images:
                image_results = self.image_parser.process_images_parallel(images)
                for result in sorted(image_results, key=lambda item: item.index):
                    if result.content:
                        parts.append(f"[Slide image content] {result.content}")

            content = "\n\n".join(parts)
            metadata = DocumentMetadata(
                parser=self.parser_name,
                document_type=DocumentType.POWERPOINT,
                size=len(file_data),
                lines=len(content.splitlines()),
            )
            return ParseResult(content, metadata)
        except Exception as exc:
            logger.error("PPTX parse failed for %s: %s", filename, exc)
            return self._empty_result(len(file_data))

    def _empty_result(self, file_size: int) -> ParseResult:
        return ParseResult(
            "",
            DocumentMetadata(
                parser=self.parser_name,
                document_type=DocumentType.POWERPOINT,
                size=file_size,
                lines=0,
            ),
        )
