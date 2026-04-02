from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ASSET_DIR = Path(r"d:\playform\police-training-platform\tmp_police_shuke_assets")
OUTPUT = Path(r"d:\QWQ\feishu\智慧教育训练平台介绍-角色终端版-警务授客风格.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

POLICE_NAVY = RGBColor(10, 36, 85)
POLICE_NAVY_DARK = RGBColor(7, 26, 64)
POLICE_BLUE = RGBColor(37, 105, 197)
POLICE_BLUE_LIGHT = RGBColor(123, 176, 255)
POLICE_GOLD = RGBColor(201, 168, 75)
POLICE_GOLD_LIGHT = RGBColor(232, 214, 161)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(245, 247, 251)
TEXT = RGBColor(48, 60, 82)
MUTED = RGBColor(115, 127, 147)
LINE = RGBColor(224, 229, 239)
RED = RGBColor(201, 64, 64)


def add_glow(base, center, radius, fill, blur):
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    x, y = center
    draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=fill)
    overlay = overlay.filter(ImageFilter.GaussianBlur(blur))
    return Image.alpha_composite(base, overlay)


def add_vertical_gradient(draw, width, height, top_color, bottom_color):
    for y in range(height):
        ratio = y / max(height - 1, 1)
        r = int(top_color[0] * (1 - ratio) + bottom_color[0] * ratio)
        g = int(top_color[1] * (1 - ratio) + bottom_color[1] * ratio)
        b = int(top_color[2] * (1 - ratio) + bottom_color[2] * ratio)
        draw.line((0, y, width, y), fill=(r, g, b, 255))


def draw_shield(draw, x, y, w, h, outline, fill=None, width=6):
    pts = [
        (x + w * 0.18, y + h * 0.10),
        (x + w * 0.82, y + h * 0.10),
        (x + w * 0.92, y + h * 0.28),
        (x + w * 0.80, y + h * 0.74),
        (x + w * 0.50, y + h * 0.96),
        (x + w * 0.20, y + h * 0.74),
        (x + w * 0.08, y + h * 0.28),
    ]
    draw.polygon(pts, outline=outline, fill=fill, width=width)


def build_cover_bg(path: Path):
    w, h = 1600, 900
    base = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(base)
    add_vertical_gradient(draw, w, h, (14, 40, 91), (10, 26, 58))

    base = add_glow(base, (260, 190), 250, (31, 112, 228, 115), 120)
    base = add_glow(base, (1230, 140), 200, (214, 175, 86, 85), 100)
    base = add_glow(base, (1200, 640), 220, (82, 135, 230, 95), 110)
    base = add_glow(base, (1350, 260), 80, (255, 88, 88, 110), 28)
    base = add_glow(base, (1450, 285), 80, (65, 148, 255, 120), 28)

    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)

    # abstract device silhouettes to mirror reference composition
    d.rounded_rectangle((320, 125, 1100, 565), radius=26, fill=(255, 255, 255, 18), outline=(255, 255, 255, 32), width=2)
    d.rounded_rectangle((365, 165, 1055, 510), radius=12, fill=(18, 60, 129, 90), outline=(146, 184, 255, 42), width=2)
    d.rectangle((555, 565, 865, 590), fill=(255, 255, 255, 28))
    d.polygon([(505, 590), (915, 590), (975, 655), (445, 655)], fill=(255, 255, 255, 20), outline=(255, 255, 255, 26))

    d.rounded_rectangle((1008, 220, 1190, 528), radius=24, fill=(255, 255, 255, 26), outline=(255, 255, 255, 36), width=2)
    d.rounded_rectangle((1026, 255, 1172, 493), radius=16, fill=(14, 83, 160, 95))
    for y in range(0, 4):
        d.rounded_rectangle((395, 195 + y * 62, 715, 238 + y * 62), radius=10, fill=(255, 255, 255, 22))
    for y in range(0, 4):
        d.rounded_rectangle((740, 195 + y * 62, 1008, 238 + y * 62), radius=10, fill=(255, 255, 255, 16))
    for y in range(0, 4):
        d.rounded_rectangle((1046, 286 + y * 48, 1148, 315 + y * 48), radius=8, fill=(255, 255, 255, 24))

    draw_shield(d, 720, 58, 150, 175, outline=(235, 221, 174, 160), width=7)
    d.ellipse((760, 95, 830, 165), fill=(235, 221, 174, 95))
    d.rectangle((770, 124, 820, 136), fill=(255, 255, 255, 180))
    d.rectangle((789, 105, 801, 155), fill=(255, 255, 255, 180))

    # bottom white wave
    d.ellipse((-240, 540, 1840, 1330), fill=(246, 248, 252, 255))
    d.ellipse((-250, 455, 1850, 1128), outline=(201, 168, 75, 255), width=18)
    d.ellipse((-258, 470, 1858, 1143), outline=(94, 153, 247, 185), width=10)

    base = Image.alpha_composite(base, overlay)
    base.convert("RGB").save(path)


def build_contents_bg(path: Path):
    w, h = 1600, 900
    base = Image.open(ASSET_DIR / "cover.png").convert("RGBA")
    veil = Image.new("RGBA", (w, h), (12, 24, 50, 110))
    base = Image.alpha_composite(base, veil)
    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    d.rectangle((0, 380, w, 900), fill=(244, 246, 251, 245))
    base = Image.alpha_composite(base, overlay)
    base.convert("RGB").save(path)


def build_section_bg(path: Path):
    w, h = 1600, 900
    base = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(base)
    add_vertical_gradient(draw, w, h, (9, 31, 74), (8, 22, 52))
    base = add_glow(base, (210, 170), 240, (33, 101, 193, 105), 110)
    base = add_glow(base, (1310, 170), 190, (204, 169, 79, 80), 100)
    base = add_glow(base, (1260, 660), 240, (69, 122, 225, 90), 120)

    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)

    draw_shield(d, 1045, 145, 360, 420, outline=(214, 194, 142, 52), width=10)
    draw_shield(d, 1080, 180, 290, 350, outline=(255, 255, 255, 28), width=5)
    for y in range(130, 700, 80):
        d.line((690, y, 1030, y - 60), fill=(90, 143, 255, 36), width=2)
        d.line((820, y, 1150, y + 74), fill=(194, 161, 81, 28), width=2)
    for x in range(250, 1450, 150):
        d.ellipse((x, 120 + (x % 120), x + 8, 128 + (x % 120)), fill=(201, 168, 75, 55))
    base = Image.alpha_composite(base, overlay)
    base.convert("RGB").save(path)


def build_light_bg(path: Path):
    w, h = 1600, 900
    base = Image.new("RGBA", (w, h), (246, 247, 251, 255))
    base = add_glow(base, (70, 820), 300, (103, 172, 255, 70), 120)
    base = add_glow(base, (1520, 90), 260, (222, 204, 154, 65), 110)
    base = add_glow(base, (1380, 660), 240, (109, 149, 255, 95), 105)

    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    d.ellipse((1240, 365, 1810, 965), fill=(255, 255, 255, 178))
    d.ellipse((1288, 412, 1760, 920), outline=(201, 168, 75, 30), width=14)
    d.ellipse((1315, 440, 1733, 892), outline=(112, 154, 255, 34), width=18)
    d.ellipse((1345, 470, 1703, 862), outline=(112, 154, 255, 24), width=12)
    draw_shield(d, 1285, 495, 310, 340, outline=(201, 168, 75, 24), width=7)
    base = Image.alpha_composite(base, overlay)
    base.convert("RGB").save(path)


def build_assets():
    ASSET_DIR.mkdir(exist_ok=True)
    build_cover_bg(ASSET_DIR / "cover.png")
    build_contents_bg(ASSET_DIR / "contents.png")
    build_section_bg(ASSET_DIR / "section.png")
    build_light_bg(ASSET_DIR / "light.png")


def add_picture_bg(slide, image_path: Path):
    slide.shapes.add_picture(str(image_path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_shadow_card(slide, left, top, width, height, fill=WHITE, transparency=0.0):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.05), top + Inches(0.06), width, height, RGBColor(203, 210, 225), transparency=0.62)
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill, line=LINE, transparency=transparency)


def add_textbox(slide, left, top, width, height, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, italic=False, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, size=15.5, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.bullet = True
        p.left_margin = Pt(18)
        p.hanging = Pt(-10)
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_logo(slide):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.56), Inches(0.24), Inches(0.44), Inches(0.44), POLICE_NAVY)
    add_textbox(slide, Inches(11.61), Inches(0.31), Inches(0.34), Inches(0.16), "警", 15, POLICE_GOLD_LIGHT, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(12.04), Inches(0.28), Inches(0.92), Inches(0.20), "智慧训练", 11, RGBColor(94, 103, 125), True)


def add_center_title(slide, title, subtitle="", module_tags=None):
    add_logo(slide)
    add_textbox(slide, Inches(2.0), Inches(0.80), Inches(9.2), Inches(0.34), title, 21, POLICE_BLUE, True, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.02), Inches(1.23), Inches(0.52), Inches(0.02), POLICE_GOLD)
    if subtitle:
        add_textbox(slide, Inches(1.08), Inches(1.42), Inches(10.9), Inches(0.42), subtitle, 13.2, MUTED, False, PP_ALIGN.CENTER)
    if module_tags:
        x = Inches(4.7)
        for tag in module_tags:
            fill = POLICE_GOLD if tag in {"管", "评"} else POLICE_BLUE
            add_pill(slide, x, Inches(1.88), Inches(0.42), tag, fill, font_size=10.5)
            x += Inches(0.52)


def add_pill(slide, left, top, width, text, fill, font_size=11):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.34), fill)
    add_textbox(slide, left + Inches(0.04), top + Inches(0.06), width - Inches(0.08), Inches(0.16), text, font_size, WHITE, True, PP_ALIGN.CENTER)


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "cover.png")
    add_textbox(slide, Inches(2.18), Inches(2.54), Inches(9.0), Inches(0.84), "智慧教育训练平台", 34, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.65), Inches(3.26), Inches(8.1), Inches(0.45), "教官 / 学员角色汇报", 22, POLICE_GOLD_LIGHT, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.72), Inches(3.95), Inches(5.95), Inches(0.24), "警务风格 · 角色 × 终端 · 双端协同汇报", 13.5, RGBColor(224, 231, 242), False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.48), Inches(6.56), Inches(2.35), Inches(0.20), "学 · 练 · 考 · 评 · 管", 11, RGBColor(85, 95, 120), False, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.18), Inches(4.58), Inches(0.08), Inches(0.08), WHITE)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.36), Inches(4.58), Inches(0.08), Inches(0.08), RGBColor(217, 221, 231))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.54), Inches(4.58), Inches(0.08), Inches(0.08), RGBColor(217, 221, 231))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.72), Inches(4.58), Inches(0.08), Inches(0.08), RGBColor(217, 221, 231))
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.62), Inches(6.96), Inches(0.42), Inches(0.42), POLICE_NAVY)
    add_textbox(slide, Inches(5.67), Inches(7.03), Inches(0.32), Inches(0.16), "警", 14, POLICE_GOLD_LIGHT, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.06), Inches(7.02), Inches(1.0), Inches(0.18), "智慧训练", 12, RGBColor(64, 74, 96), True)
    return slide


def contents_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "contents.png")
    add_textbox(slide, Inches(4.14), Inches(1.22), Inches(5.1), Inches(0.60), "CONTENTS", 34, POLICE_BLUE_LIGHT, False, PP_ALIGN.CENTER, True, "Calibri")
    cards = [
        ("01", "双端定位", ["角色 × 终端", "汇报逻辑"]),
        ("02", "五模块总览", ["学练考评管", "能力映射"]),
        ("03", "教官端能力", ["网页端管理", "移动端带班"]),
        ("04", "学员端能力", ["网页端学习考试", "移动端参训"]),
        ("05", "双端协同", ["典型场景", "成果与规划"]),
    ]
    start = Inches(0.62)
    for idx, (num, title, bullets) in enumerate(cards):
        x = start + Inches(idx * 2.50)
        add_shadow_card(slide, x, Inches(3.28), Inches(2.02), Inches(2.14), fill=WHITE, transparency=0.14)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.28), Inches(2.02), Inches(0.54), WHITE, transparency=0.36)
        add_textbox(slide, x + Inches(0.54), Inches(3.51), Inches(0.94), Inches(0.30), num, 22, WHITE, False, PP_ALIGN.CENTER, True, "Calibri")
        add_textbox(slide, x + Inches(0.11), Inches(4.16), Inches(1.80), Inches(0.28), title, 15.2, TEXT, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.20), Inches(4.60), Inches(1.62), Inches(0.56), "\n".join([f"· {item}" for item in bullets]), 11.5, MUTED)
    return slide


def positioning_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(
        slide,
        "同一套训练业务的双端定位",
        "网页端负责管理与深度操作，移动端负责现场执行与高频使用。",
        module_tags=["学", "练", "考", "评", "管"],
    )
    cards = [
        ("教官网页端", "培训班组织\n课程与资源\n考试协同", POLICE_GOLD),
        ("教官移动端", "班级查看\n签到签退\n现场通知", POLICE_BLUE),
        ("学员网页端", "在线学习\n考试参与\n训练记录", POLICE_GOLD),
        ("学员移动端", "我的班级\n扫码签到\n移动学习", POLICE_BLUE),
    ]
    xs = [Inches(0.95), Inches(3.95), Inches(6.95), Inches(9.95)]
    for (title, body, fill), x in zip(cards, xs):
        add_shadow_card(slide, x, Inches(2.40), Inches(2.35), Inches(3.15))
        add_pill(slide, x + Inches(0.34), Inches(2.22), Inches(1.68), title, fill)
        add_textbox(slide, x + Inches(0.18), Inches(2.96), Inches(2.0), Inches(1.55), body, 15, TEXT, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.55), Inches(6.24), Inches(10.2), Inches(0.32), "后续汇报围绕角色 × 终端展开，不再把两端拆成两套系统，而是说明同一条训练链路在不同场景中的使用方式。", 13.2, MUTED, False, PP_ALIGN.CENTER)
    return slide


def module_map_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(slide, "学练考评管五模块映射", "把五个模块映射到角色与终端中，更清晰说明系统能力边界。", module_tags=["学", "练", "考", "评", "管"])
    cards = [
        ("学", "课程资源\n资源库\n移动学习", "对应网页端课程资源、资源库，以及移动端课程查看与学习。"),
        ("练", "培训班\n排课\n签到签退", "对应网页端培训班管理与移动端班级执行、现场签到。"),
        ("考", "题库协同\n在线考试\n结果查看", "对应考试安排、考试参与、结果回传与班级考试协同。"),
        ("评", "训练记录\n个人训历\n结果沉淀", "对应学习、签到、考试结果沉淀形成个人训练画像。"),
        ("管", "工作台\n角色权限\n统一支撑", "对应网页端组织管理、权限治理和双端统一支撑。"),
    ]
    xs = [Inches(0.62), Inches(3.17), Inches(5.72), Inches(8.27), Inches(10.82)]
    fills = [POLICE_BLUE, POLICE_GOLD, POLICE_BLUE, POLICE_GOLD, POLICE_BLUE]
    for (tag, title, desc), x, fill in zip(cards, xs, fills):
        add_shadow_card(slide, x, Inches(2.18), Inches(2.0), Inches(3.52))
        add_pill(slide, x + Inches(0.58), Inches(2.00), Inches(0.84), tag, fill, 12)
        add_textbox(slide, x + Inches(0.16), Inches(2.72), Inches(1.68), Inches(0.80), title, 15, TEXT, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.14), Inches(3.56), Inches(1.70), Inches(1.40), desc, 11.5, MUTED, False, PP_ALIGN.CENTER)
    return slide


def loop_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(slide, "角色与终端协同下的核心业务闭环", "围绕训练组织、学习、签到、考试和结果沉淀形成一体化闭环。", module_tags=["学", "练", "考", "评", "管"])
    steps = [
        "人员进入系统",
        "培训班组织",
        "课程与资源",
        "现场签到",
        "在线考试",
        "训练记录",
        "结果分析",
    ]
    widths = [1.56, 1.62, 1.58, 1.45, 1.45, 1.55, 1.43]
    fills = [POLICE_GOLD, POLICE_BLUE, POLICE_GOLD, POLICE_BLUE, POLICE_GOLD, POLICE_BLUE, POLICE_GOLD]
    x = Inches(0.68)
    for idx, (step, width, fill) in enumerate(zip(steps, widths, fills)):
        add_shadow_card(slide, x, Inches(3.10), Inches(width), Inches(1.24))
        add_pill(slide, x + Inches(0.18), Inches(2.92), Inches(0.68), f"0{idx + 1}", fill)
        add_textbox(slide, x + Inches(0.10), Inches(3.58), Inches(width) - Inches(0.20), Inches(0.36), step, 13.2, TEXT, True, PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_textbox(slide, x + Inches(width) + Inches(0.02), Inches(3.51), Inches(0.26), Inches(0.24), "→", 23, POLICE_BLUE, True, PP_ALIGN.CENTER, font="Calibri")
        x += Inches(width) + Inches(0.17)
    add_shadow_card(slide, Inches(1.12), Inches(5.30), Inches(4.92), Inches(0.98))
    add_pill(slide, Inches(1.36), Inches(5.13), Inches(1.60), "网页端承担", POLICE_BLUE)
    add_textbox(slide, Inches(1.42), Inches(5.68), Inches(4.30), Inches(0.24), "组织、配置、查看、分析", 16, TEXT, True, PP_ALIGN.CENTER)
    add_shadow_card(slide, Inches(7.18), Inches(5.30), Inches(4.92), Inches(0.98))
    add_pill(slide, Inches(7.42), Inches(5.13), Inches(1.60), "移动端承担", POLICE_GOLD)
    add_textbox(slide, Inches(7.48), Inches(5.68), Inches(4.30), Inches(0.24), "签到、查看、执行、提醒", 16, TEXT, True, PP_ALIGN.CENTER)
    return slide


def section_slide(prs, english, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "section.png")
    add_textbox(slide, Inches(4.28), Inches(1.82), Inches(4.8), Inches(0.48), english, 24, POLICE_GOLD_LIGHT, False, PP_ALIGN.CENTER, True, "Calibri")
    add_textbox(slide, Inches(2.14), Inches(2.98), Inches(9.0), Inches(0.62), title, 28, WHITE, True, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.02), Inches(3.78), Inches(0.54), Inches(0.02), POLICE_GOLD)
    add_textbox(slide, Inches(1.82), Inches(4.16), Inches(9.7), Inches(0.62), subtitle, 16, RGBColor(223, 228, 239), False, PP_ALIGN.CENTER)
    return slide


def dual_cards_slide(prs, title, subtitle, left_title, left_items, right_title, right_items, tags):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(slide, title, subtitle, module_tags=tags)
    add_shadow_card(slide, Inches(0.85), Inches(2.15), Inches(5.74), Inches(4.12))
    add_pill(slide, Inches(1.12), Inches(1.97), Inches(1.95), left_title, POLICE_GOLD)
    add_bullets(slide, Inches(1.06), Inches(2.68), Inches(5.02), Inches(3.16), left_items)
    add_shadow_card(slide, Inches(6.74), Inches(2.15), Inches(5.74), Inches(4.12))
    add_pill(slide, Inches(7.01), Inches(1.97), Inches(1.95), right_title, POLICE_BLUE)
    add_bullets(slide, Inches(6.95), Inches(2.68), Inches(5.02), Inches(3.16), right_items)
    return slide


def four_card_slide(prs, title, subtitle, cards, tags):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(slide, title, subtitle, module_tags=tags)
    positions = [
        (Inches(0.92), Inches(2.18)),
        (Inches(3.97), Inches(2.18)),
        (Inches(7.02), Inches(2.18)),
        (Inches(10.07), Inches(2.18)),
    ]
    fills = [POLICE_GOLD, POLICE_BLUE, POLICE_GOLD, POLICE_BLUE]
    for (header, body), (left, top), fill in zip(cards, positions, fills):
        add_shadow_card(slide, left, top, Inches(2.35), Inches(3.60))
        add_pill(slide, left + Inches(0.36), top - Inches(0.15), Inches(1.56), header, fill)
        add_textbox(slide, left + Inches(0.16), top + Inches(0.78), Inches(2.02), Inches(2.34), body, 14.5, TEXT, False, PP_ALIGN.CENTER)
    return slide


def scene_list_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(slide, "双端协同典型场景", "不是两套系统各讲各的，而是同一条训练链路在不同终端间自然衔接。", module_tags=["学", "练", "考", "评", "管"])
    items = [
        ("培训班组织", "教官网页端完成班级组织、课表安排和资源准备，学员通过网页端或移动端查看班级与日历。", POLICE_GOLD),
        ("现场签到", "教官移动端发起签到签退，学员移动端扫码完成签到，结果实时回流到班级记录。", POLICE_BLUE),
        ("学习支撑", "教官网页端准备课程资源与教学内容，学员网页端和移动端持续查看课程与资源。", POLICE_GOLD),
        ("考试闭环", "教官网页端协同考试安排，学员网页端或移动端参加考试，成绩统一沉淀到训练记录。", POLICE_BLUE),
    ]
    top = Inches(2.00)
    for idx, (head, desc, fill) in enumerate(items):
        y = top + Inches(idx * 1.08)
        add_shadow_card(slide, Inches(1.06), y, Inches(11.0), Inches(0.82))
        add_pill(slide, Inches(1.28), y + Inches(0.20), Inches(1.34), head, fill)
        add_textbox(slide, Inches(2.95), y + Inches(0.21), Inches(8.68), Inches(0.34), desc, 13.0, MUTED)
    return slide


def support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(slide, "统一支撑关系", "统一角色、统一权限、统一训练数据，支撑网页端和移动端的协同运行。", module_tags=["管"])
    add_shadow_card(slide, Inches(1.02), Inches(1.88), Inches(11.25), Inches(1.0))
    add_textbox(slide, Inches(1.30), Inches(2.22), Inches(10.68), Inches(0.28), "同一套后端与数据体系，支撑教官和学员在网页端、移动端的不同使用方式。", 17, TEXT, True, PP_ALIGN.CENTER)
    center_modules = [("训练组织", POLICE_GOLD), ("课程资源", POLICE_BLUE), ("在线考试", POLICE_GOLD), ("训练记录", POLICE_BLUE)]
    x = Inches(1.58)
    for text, fill in center_modules:
        add_shadow_card(slide, x, Inches(3.10), Inches(2.28), Inches(0.94))
        add_pill(slide, x + Inches(0.34), Inches(2.92), Inches(1.60), text, fill)
        x += Inches(2.55)
    ends = [
        ("教官网页端", POLICE_GOLD, Inches(1.18)),
        ("教官移动端", POLICE_BLUE, Inches(4.08)),
        ("学员网页端", POLICE_GOLD, Inches(6.98)),
        ("学员移动端", POLICE_BLUE, Inches(9.88)),
    ]
    for text, fill, left in ends:
        add_shadow_card(slide, left, Inches(4.55), Inches(2.26), Inches(1.24))
        add_pill(slide, left + Inches(0.35), Inches(4.37), Inches(1.56), text, fill)
    return slide


def result_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(slide, "当前成果与演示重点", "聚焦当前项目中已经具备展示价值的能力，突出双端落地场景。", module_tags=["学", "练", "考", "评", "管"])
    add_shadow_card(slide, Inches(0.88), Inches(2.10), Inches(5.78), Inches(4.18))
    add_pill(slide, Inches(1.16), Inches(1.92), Inches(1.76), "网页端成果", POLICE_BLUE)
    add_bullets(
        slide,
        Inches(1.08),
        Inches(2.64),
        Inches(5.04),
        Inches(3.20),
        [
            "教官网页端已覆盖工作台、培训班管理、排课查看、课程资源、资源库、考试协同等核心能力。",
            "学员网页端已覆盖学习、考试、报名、训练历史和个人中心等关键链路。",
            "适合突出“管理与深度操作”的桌面端价值。",
        ],
    )
    add_shadow_card(slide, Inches(6.78), Inches(2.10), Inches(5.70), Inches(4.18))
    add_pill(slide, Inches(7.06), Inches(1.92), Inches(1.76), "移动端成果", POLICE_GOLD)
    add_bullets(
        slide,
        Inches(6.98),
        Inches(2.64),
        Inches(4.98),
        Inches(3.20),
        [
            "移动端已覆盖首页工作台、我的班级、日历课表、签到签退、资源查看、在线考试和通知入口。",
            "适合突出“带班现场、扫码签到、移动参训”的高频执行价值。",
            "更容易体现项目在真实训练场景中的落地性。",
        ],
    )
    return slide


def next_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, ASSET_DIR / "light.png")
    add_center_title(slide, "下一步优化建议", "在当前角色 × 终端结构上，继续增强训练闭环的完整性和汇报表达的清晰度。", module_tags=["学", "练", "考", "评", "管"])
    items = [
        ("继续补强角色页细节", "完善报名、通知、个人中心等角色端细节体验。", POLICE_GOLD),
        ("强化双端消息联动", "围绕公告、考试提醒、签到异常和待办做统一触达。", POLICE_BLUE),
        ("提升结果沉淀能力", "把签到、学习、考试进一步沉淀到个人训历与评估分析。", POLICE_GOLD),
        ("统一汇报话术", "始终围绕角色 × 终端 × 协同闭环展开，不再回到页面堆砌式表达。", POLICE_BLUE),
    ]
    positions = [
        (Inches(1.0), Inches(2.15)),
        (Inches(6.92), Inches(2.15)),
        (Inches(1.0), Inches(4.38)),
        (Inches(6.92), Inches(4.38)),
    ]
    for (head, desc, fill), (left, top) in zip(items, positions):
        add_shadow_card(slide, left, top, Inches(5.24), Inches(1.58))
        add_pill(slide, left + Inches(0.25), top + Inches(0.20), Inches(1.98), head, fill)
        add_textbox(slide, left + Inches(0.27), top + Inches(0.82), Inches(4.62), Inches(0.42), desc, 13.0, MUTED)
    return slide


def build_deck():
    build_assets()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_slide(prs)
    contents_slide(prs)
    positioning_slide(prs)
    module_map_slide(prs)
    loop_slide(prs)
    section_slide(prs, "INSTRUCTOR", "教官角色业务全景", "网页端负责组织与管理，移动端负责带班与现场执行。")
    dual_cards_slide(
        prs,
        "教官端能力说明",
        "网页端强调管理深度，移动端强调带班效率。",
        "教官网页端",
        [
            "工作台：查看近期培训班、课程和考试待办。",
            "培训班管理：查看班级详情、学员名单、考试安排和变更记录。",
            "排课与周训练计划：查看培训班课表和训练计划。",
            "教学资源：课程资源、资源库、我的资源、教学资源生成。",
            "考试协同：查看考试安排和班级考试结果。",
        ],
        "教官移动端",
        [
            "首页工作台：查看个人信息、统计卡片和近期安排。",
            "我的班级：快速进入班级列表和班级详情。",
            "日历课表：随时查看课程和考试安排。",
            "签到签退：发起签到、结束签到、发起签退。",
            "班级公告与通知：现场查看并处理通知提醒。",
        ],
        tags=["练", "考", "管", "学"],
    )
    four_card_slide(
        prs,
        "教官端汇报重点",
        "把教官端讲成“组织管理 + 现场带班”的一体化入口。",
        [
            ("组织管理", "适合讲培训班组织、课程安排、资源准备和考试协同。"),
            ("深度查看", "适合讲班级详情、学员名单、变更记录和训练计划。"),
            ("现场带班", "适合讲签到签退、班级公告、课表查看和现场执行。"),
            ("双端协同", "强调网页端负责准备，移动端负责执行，形成完整闭环。"),
        ],
        tags=["练", "考", "管"],
    )
    section_slide(prs, "STUDENT", "学员角色业务全景", "网页端负责学习与考试，移动端负责参训与现场执行。")
    dual_cards_slide(
        prs,
        "学员端能力说明",
        "网页端强调学习考试，移动端强调参训签到。",
        "学员网页端",
        [
            "学员工作台：查看学习统计、考试次数和近期安排。",
            "培训班参与：查看班级、进入详情、发起报名申请。",
            "在线学习：查看课程资源、课程详情和学习进度。",
            "在线考试：查看考试列表、考试概览、在线作答和考试结果。",
            "训练历史：查看个人参训记录和考试记录。",
        ],
        "学员移动端",
        [
            "首页：查看个人信息、统计卡片和快捷入口。",
            "我的班级：查看参加中的班级和班级详情。",
            "日历课表：按周查看训练与考试安排。",
            "签到签退：扫码签到、扫码签退或直接签到。",
            "移动学习与移动考试：随时查看资源、参加考试和接收提醒。",
        ],
        tags=["学", "练", "考", "评"],
    )
    four_card_slide(
        prs,
        "学员端汇报重点",
        "把学员端讲成“学习入口 + 参训入口”的协同体系。",
        [
            ("学习入口", "网页端更适合做连续性学习和完整考试操作。"),
            ("参训入口", "移动端更适合查看班级、扫码签到和现场执行。"),
            ("记录沉淀", "学习、考试和签到最终都会沉淀到个人训练记录中。"),
            ("体验价值", "让学员在不同场景下都能顺畅完成参训和学习。"),
        ],
        tags=["学", "练", "考", "评"],
    )
    scene_list_slide(prs)
    support_slide(prs)
    result_slide(prs)
    next_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build_deck()
