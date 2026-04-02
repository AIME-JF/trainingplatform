from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(11, 35, 84)
BLUE = RGBColor(34, 91, 184)
LIGHT_BLUE = RGBColor(232, 240, 255)
GOLD = RGBColor(201, 168, 75)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(247, 249, 252)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(99, 115, 129)
LINE = RGBColor(221, 228, 237)
GREEN = RGBColor(31, 147, 108)
ORANGE = RGBColor(229, 131, 54)
RED = RGBColor(205, 72, 72)


def add_full_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill, line=None, radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius_shape, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Microsoft YaHei",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    font_size=18,
    color=TEXT,
    level_indent=18,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.left_margin = Pt(level_indent)
        p.hanging = Pt(-12)
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return box


def add_title_band(slide, title, subtitle="", page_text=""):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.72), NAVY, radius_shape=MSO_SHAPE.RECTANGLE)
    add_rect(slide, Inches(0.44), Inches(0.21), Inches(0.12), Inches(0.30), GOLD, radius_shape=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(0.68), Inches(0.15), Inches(6.7), Inches(0.34), title, 26, WHITE, True)
    if subtitle:
        add_textbox(slide, Inches(0.70), Inches(0.42), Inches(8.5), Inches(0.20), subtitle, 10.5, RGBColor(214, 224, 240))
    if page_text:
        add_textbox(slide, Inches(11.85), Inches(0.18), Inches(1.0), Inches(0.26), page_text, 10, WHITE, False, PP_ALIGN.RIGHT)


def add_footer(slide, text="智慧教育训练平台 · 教官/学员双端汇报"):
    add_textbox(slide, Inches(0.6), Inches(7.0), Inches(5.5), Inches(0.2), text, 9.5, MUTED)


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, NAVY)
    add_rect(slide, Inches(0.65), Inches(0.72), Inches(0.34), Inches(5.75), GOLD, radius_shape=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(1.25), Inches(1.08), Inches(5.6), Inches(0.6), "智慧教育训练平台", 28, WHITE, True)
    add_textbox(slide, Inches(1.25), Inches(1.82), Inches(6.7), Inches(0.8), "教官 / 学员角色汇报", 24, WHITE, True)
    add_textbox(slide, Inches(1.25), Inches(2.48), Inches(7.6), Inches(0.52), "网页端 + 移动端协同视角", 20, RGBColor(219, 229, 241), False)
    add_textbox(
        slide,
        Inches(1.25),
        Inches(3.18),
        Inches(8.0),
        Inches(1.2),
        "围绕“角色 × 终端”展开，突出同一套训练业务在不同场景中的使用方式，\n体现教官网页端管理、教官移动端带班、学员网页端学习、学员移动端参训的完整闭环。",
        16,
        RGBColor(228, 234, 244),
    )

    tags = [
        ("同一套训练业务", BLUE),
        ("网页端负责管理与深度操作", RGBColor(25, 104, 185)),
        ("移动端负责现场执行与高频使用", RGBColor(21, 137, 101)),
    ]
    left = Inches(1.25)
    top = Inches(5.05)
    for text, color in tags:
        w = Inches(2.2 if len(text) < 9 else 3.05 if len(text) < 15 else 4.15)
        add_rect(slide, left, top, w, Inches(0.52), color)
        add_textbox(slide, left + Inches(0.14), top + Inches(0.11), w - Inches(0.28), Inches(0.24), text, 12, WHITE, True, PP_ALIGN.CENTER)
        left += w + Inches(0.18)

    add_rect(slide, Inches(9.15), Inches(1.0), Inches(3.15), Inches(4.9), RGBColor(19, 47, 102))
    quad_titles = [
        ("教官网页端", "培训组织 / 资源管理 / 考试协同"),
        ("教官移动端", "班级现场 / 签到签退 / 公告通知"),
        ("学员网页端", "在线学习 / 考试参与 / 训练记录"),
        ("学员移动端", "我的班级 / 扫码签到 / 移动学习"),
    ]
    positions = [
        (Inches(9.45), Inches(1.35)),
        (Inches(10.92), Inches(1.35)),
        (Inches(9.45), Inches(3.55)),
        (Inches(10.92), Inches(3.55)),
    ]
    for (title, desc), (left, top) in zip(quad_titles, positions):
        add_rect(slide, left, top, Inches(1.18), Inches(1.75), WHITE)
        add_textbox(slide, left + Inches(0.08), top + Inches(0.20), Inches(1.02), Inches(0.45), title, 12.5, NAVY, True, PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.08), top + Inches(0.72), Inches(1.02), Inches(0.70), desc, 9.5, MUTED, False, PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.25), Inches(6.56), Inches(4.8), Inches(0.24), "2026年4月", 11, RGBColor(214, 224, 240))
    return slide


def agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, "汇报目录", "基于角色 × 终端的双端能力汇报", "2 / 14")
    items = [
        ("01", "建设目标与双端定位", "说明为什么要按“角色 × 终端”汇报"),
        ("02", "教官端能力", "网页端管理、移动端现场带班"),
        ("03", "学员端能力", "网页端学习考试、移动端参训执行"),
        ("04", "双端协同与项目价值", "强调不是两套系统，而是一套业务闭环"),
        ("05", "当前成果与下一步", "明确已实现范围和后续优化方向"),
    ]
    top = Inches(1.2)
    for idx, (num, title, desc) in enumerate(items):
        add_rect(slide, Inches(0.8), top + Inches(idx * 1.0), Inches(11.7), Inches(0.78), WHITE, LINE)
        add_rect(slide, Inches(1.0), top + Inches(idx * 1.0) + Inches(0.15), Inches(0.85), Inches(0.48), NAVY)
        add_textbox(slide, Inches(1.02), top + Inches(idx * 1.0) + Inches(0.23), Inches(0.80), Inches(0.18), num, 16, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.05), top + Inches(idx * 1.0) + Inches(0.12), Inches(4.2), Inches(0.25), title, 19, TEXT, True)
        add_textbox(slide, Inches(2.05), top + Inches(idx * 1.0) + Inches(0.40), Inches(7.8), Inches(0.20), desc, 11.5, MUTED)
    add_footer(slide)
    return slide


def positioning_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, "双端定位", "同一套训练业务，网页端负责管理与深度操作，移动端负责现场执行与高频使用", "3 / 14")
    add_rect(slide, Inches(0.75), Inches(1.1), Inches(5.95), Inches(4.95), WHITE, LINE)
    add_textbox(slide, Inches(1.0), Inches(1.34), Inches(2.0), Inches(0.3), "网页端定位", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(1.8),
        Inches(5.2),
        Inches(3.5),
        [
            "适合培训组织、班级配置、课程资源、考试协同等信息量大、流程长的操作。",
            "强调全局视角、数据查看、深度管理和完整配置能力。",
            "更适合作为教官和学员的桌面学习与管理入口。",
        ],
        16,
    )
    add_rect(slide, Inches(6.85), Inches(1.1), Inches(5.75), Inches(4.95), WHITE, LINE)
    add_textbox(slide, Inches(7.1), Inches(1.34), Inches(2.0), Inches(0.3), "移动端定位", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(7.1),
        Inches(1.8),
        Inches(5.0),
        Inches(3.5),
        [
            "适合班级现场执行、签到签退、课表查看、通知提醒等高频轻操作。",
            "强调即看即用、随身可达、培训现场可执行。",
            "更适合作为教官带班和学员参训的日常入口。",
        ],
        16,
    )
    add_rect(slide, Inches(1.0), Inches(5.55), Inches(2.6), Inches(0.8), LIGHT_BLUE, BLUE)
    add_textbox(slide, Inches(1.15), Inches(5.72), Inches(2.25), Inches(0.36), "教官网页端：组织管理", 16, NAVY, True, PP_ALIGN.CENTER)
    add_rect(slide, Inches(3.95), Inches(5.55), Inches(2.6), Inches(0.8), LIGHT_BLUE, BLUE)
    add_textbox(slide, Inches(4.10), Inches(5.72), Inches(2.25), Inches(0.36), "教官移动端：现场带班", 16, NAVY, True, PP_ALIGN.CENTER)
    add_rect(slide, Inches(6.90), Inches(5.55), Inches(2.6), Inches(0.8), LIGHT_BLUE, BLUE)
    add_textbox(slide, Inches(7.05), Inches(5.72), Inches(2.25), Inches(0.36), "学员网页端：学习考试", 16, NAVY, True, PP_ALIGN.CENTER)
    add_rect(slide, Inches(9.85), Inches(5.55), Inches(2.6), Inches(0.8), LIGHT_BLUE, BLUE)
    add_textbox(slide, Inches(10.00), Inches(5.72), Inches(2.25), Inches(0.36), "学员移动端：参训执行", 16, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide)
    return slide


def loop_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, "核心业务闭环", "通过角色协同与双端协同，把训练组织、学习、考试和结果沉淀串成完整链路", "4 / 14")
    steps = [
        "人员与角色进入系统",
        "培训班组织与排课",
        "课程学习与资源查看",
        "现场签到与执行反馈",
        "在线考试与结果回传",
        "训练历史与个人档案",
        "统计分析与持续优化",
    ]
    left = Inches(0.78)
    for idx, step in enumerate(steps):
        w = Inches(1.62) if idx not in {1, 5} else Inches(1.78)
        add_rect(slide, left, Inches(2.45), w, Inches(1.15), WHITE, LINE)
        add_rect(slide, left + Inches(0.10), Inches(2.58), Inches(0.42), Inches(0.32), NAVY, radius_shape=MSO_SHAPE.RECTANGLE)
        add_textbox(slide, left + Inches(0.11), Inches(2.63), Inches(0.38), Inches(0.18), f"{idx + 1}", 12, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.16), Inches(3.02), w - Inches(0.32), Inches(0.38), step, 14, TEXT, True, PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_textbox(slide, left + w + Inches(0.03), Inches(2.90), Inches(0.28), Inches(0.24), "→", 24, BLUE, True, PP_ALIGN.CENTER)
        left += w + Inches(0.20)
    add_rect(slide, Inches(1.05), Inches(4.5), Inches(5.2), Inches(1.18), RGBColor(235, 243, 255), BLUE)
    add_textbox(slide, Inches(1.22), Inches(4.72), Inches(4.85), Inches(0.52), "网页端重点承担：组织、配置、查看、分析", 18, NAVY, True, PP_ALIGN.CENTER)
    add_rect(slide, Inches(7.05), Inches(4.5), Inches(5.2), Inches(1.18), RGBColor(235, 248, 242), GREEN)
    add_textbox(slide, Inches(7.22), Inches(4.72), Inches(4.85), Inches(0.52), "移动端重点承担：签到、查看、执行、提醒", 18, GREEN, True, PP_ALIGN.CENTER)
    add_footer(slide)
    return slide


def role_overview_slide(prs, title, subtitle, page_text, left_title, left_points, right_title, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, title, subtitle, page_text)
    add_rect(slide, Inches(0.8), Inches(1.2), Inches(12.0), Inches(0.78), NAVY)
    add_textbox(slide, Inches(1.0), Inches(1.42), Inches(11.5), Inches(0.24), "核心目标：通过双端协同，把同一角色在不同场景下的工作链路跑通，而不是把功能拆散理解。", 16, WHITE, True)
    add_rect(slide, Inches(0.85), Inches(2.2), Inches(5.75), Inches(3.8), WHITE, LINE)
    add_textbox(slide, Inches(1.08), Inches(2.45), Inches(4.0), Inches(0.25), left_title, 20, NAVY, True)
    add_bullets(slide, Inches(1.05), Inches(2.9), Inches(5.1), Inches(2.9), left_points, 16)
    add_rect(slide, Inches(6.75), Inches(2.2), Inches(5.75), Inches(3.8), WHITE, LINE)
    add_textbox(slide, Inches(6.98), Inches(2.45), Inches(4.0), Inches(0.25), right_title, 20, NAVY, True)
    add_bullets(slide, Inches(6.95), Inches(2.9), Inches(5.1), Inches(2.9), right_points, 16)
    add_footer(slide)
    return slide


def capability_slide(prs, title, subtitle, page_text, left_header, left_items, right_header, right_items, accent=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, title, subtitle, page_text)
    add_rect(slide, Inches(0.8), Inches(1.15), Inches(5.85), Inches(5.65), WHITE, LINE)
    add_rect(slide, Inches(0.8), Inches(1.15), Inches(5.85), Inches(0.68), accent, radius_shape=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(1.02), Inches(1.34), Inches(2.8), Inches(0.26), left_header, 20, WHITE, True)
    add_bullets(slide, Inches(1.0), Inches(1.98), Inches(5.1), Inches(4.5), left_items, 16)
    add_rect(slide, Inches(6.75), Inches(1.15), Inches(5.75), Inches(5.65), WHITE, LINE)
    add_rect(slide, Inches(6.75), Inches(1.15), Inches(5.75), Inches(0.68), NAVY, radius_shape=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(6.98), Inches(1.34), Inches(2.8), Inches(0.26), right_header, 20, WHITE, True)
    add_bullets(slide, Inches(6.95), Inches(1.98), Inches(5.0), Inches(4.5), right_items, 16)
    add_footer(slide)
    return slide


def scene_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, "双端协同典型场景", "突出“同一套训练业务”如何在不同终端间流转和回流", "11 / 14")
    scenes = [
        ("场景一", "培训班组织", "教官网页端完成班级组织、课表安排和资源准备，学员通过网页端或移动端查看班级与日历。"),
        ("场景二", "现场签到", "教官在移动端发起签到签退，学员通过移动端扫码完成签到，结果实时回流到班级记录。"),
        ("场景三", "学习支撑", "教官网页端准备课程资源与教学内容，学员在网页端或移动端进行在线学习和资源查看。"),
        ("场景四", "考试闭环", "教官网页端配置考试安排，学员在网页端或移动端参与考试，成绩与训练记录统一沉淀。"),
    ]
    positions = [
        (Inches(0.85), Inches(1.35)),
        (Inches(6.75), Inches(1.35)),
        (Inches(0.85), Inches(4.0)),
        (Inches(6.75), Inches(4.0)),
    ]
    colors = [BLUE, GREEN, ORANGE, RED]
    for (tag, title, desc), (left, top), color in zip(scenes, positions, colors):
        add_rect(slide, left, top, Inches(5.7), Inches(2.15), WHITE, LINE)
        add_rect(slide, left + Inches(0.18), top + Inches(0.18), Inches(0.85), Inches(0.38), color, radius_shape=MSO_SHAPE.RECTANGLE)
        add_textbox(slide, left + Inches(0.20), top + Inches(0.24), Inches(0.80), Inches(0.18), tag, 11.5, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(1.18), top + Inches(0.18), Inches(2.6), Inches(0.25), title, 18, TEXT, True)
        add_textbox(slide, left + Inches(0.20), top + Inches(0.72), Inches(5.15), Inches(1.0), desc, 14.5, MUTED)
    add_footer(slide)
    return slide


def support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, "统一支撑关系", "不是两套独立系统，而是统一角色、统一权限、统一数据和统一业务流程", "12 / 14")
    add_rect(slide, Inches(1.0), Inches(1.35), Inches(11.3), Inches(1.0), WHITE, LINE)
    add_textbox(slide, Inches(1.25), Inches(1.68), Inches(10.8), Inches(0.35), "同一套后端与数据体系，支撑教官和学员在网页端、移动端的不同使用方式。", 18, NAVY, True, PP_ALIGN.CENTER)
    modules = [
        "训练组织",
        "课程资源",
        "在线考试",
        "通知公告",
        "个人档案",
        "AI辅助",
    ]
    left = Inches(1.1)
    for module in modules:
        add_rect(slide, left, Inches(3.0), Inches(1.55), Inches(0.9), LIGHT_BLUE, BLUE)
        add_textbox(slide, left + Inches(0.05), Inches(3.30), Inches(1.45), Inches(0.20), module, 14.5, NAVY, True, PP_ALIGN.CENTER)
        left += Inches(1.82)
    add_textbox(slide, Inches(1.6), Inches(4.35), Inches(10.1), Inches(0.35), "教官网页端  ←→  教官移动端  ←→  学员网页端  ←→  学员移动端", 20, BLUE, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.9), Inches(5.0), Inches(9.5), Inches(0.65), "前端分工不同，但最终都会回到同一套训练班级、课程、考试、签到、通知和训练记录中。", 17, MUTED, False, PP_ALIGN.CENTER)
    add_footer(slide)
    return slide


def result_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, "当前成果与演示重点", "聚焦当前项目中已经具备展示价值的能力，避免把规划能力与现状混讲", "13 / 14")
    add_rect(slide, Inches(0.8), Inches(1.2), Inches(5.8), Inches(5.3), WHITE, LINE)
    add_textbox(slide, Inches(1.05), Inches(1.45), Inches(2.8), Inches(0.25), "网页端当前重点能力", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(1.92),
        Inches(5.0),
        Inches(4.2),
        [
            "教官网页端已覆盖工作台、培训班管理、排课查看、课程资源、资源库、考试协同等核心功能。",
            "学员网页端已覆盖在线学习、考试参与、报名申请、训练历史与个人中心等功能。",
            "适合演示“深度操作、完整配置、数据查看”的桌面使用场景。",
        ],
        16,
    )
    add_rect(slide, Inches(6.8), Inches(1.2), Inches(5.7), Inches(5.3), WHITE, LINE)
    add_textbox(slide, Inches(7.05), Inches(1.45), Inches(2.8), Inches(0.25), "移动端当前重点能力", 20, NAVY, True)
    add_bullets(
        slide,
        Inches(7.0),
        Inches(1.92),
        Inches(4.95),
        Inches(4.2),
        [
            "移动端已覆盖首页工作台、我的班级、日历课表、签到签退、资源查看、在线考试和通知入口。",
            "适合演示“带班现场、扫码签到、碎片化学习、即时提醒”的高频场景。",
            "更能体现项目在实际训练现场的执行价值和使用便捷性。",
        ],
        16,
    )
    add_footer(slide)
    return slide


def next_step_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide, BG)
    add_title_band(slide, "下一步优化建议", "在现有角色 × 终端结构上，继续增强训练闭环的完整性和落地性", "14 / 14")
    items = [
        ("01", "完善角色页细节", "继续补强学员报名、个人中心、通知中心等角色端细节体验。"),
        ("02", "强化双端消息联动", "围绕班级公告、签到异常、考试提醒和待办消息做统一触达。"),
        ("03", "打通训练结果沉淀", "把签到、学习、考试进一步汇总到个人训历和后续评估分析中。"),
        ("04", "提升现场执行能力", "继续优化教官带班操作链路和学员参训操作链路，让移动端更贴近现场。"),
        ("05", "统一汇报口径", "后续汇报始终按照‘角色 × 终端 × 协同闭环’展开，避免回到页面堆砌式表达。"),
    ]
    top = Inches(1.2)
    for idx, (num, title, desc) in enumerate(items):
        add_rect(slide, Inches(0.9), top + Inches(idx * 0.98), Inches(11.6), Inches(0.72), WHITE, LINE)
        add_rect(slide, Inches(1.08), top + Inches(idx * 0.98) + Inches(0.14), Inches(0.72), Inches(0.42), NAVY, radius_shape=MSO_SHAPE.RECTANGLE)
        add_textbox(slide, Inches(1.09), top + Inches(idx * 0.98) + Inches(0.22), Inches(0.68), Inches(0.18), num, 13, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.98), top + Inches(idx * 0.98) + Inches(0.12), Inches(3.0), Inches(0.22), title, 18, TEXT, True)
        add_textbox(slide, Inches(5.05), top + Inches(idx * 0.98) + Inches(0.13), Inches(6.2), Inches(0.30), desc, 12.5, MUTED)
    add_footer(slide)
    return slide


def build_deck(output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_slide(prs)
    agenda_slide(prs)
    positioning_slide(prs)
    loop_slide(prs)
    role_overview_slide(
        prs,
        "教官角色业务全景",
        "教官的核心目标：把训练组织起来、执行到位、结果留痕",
        "5 / 14",
        "教官网页端价值",
        [
            "承担培训班管理、排课查看、课程资源准备、考试协同等深度操作。",
            "适合完整查看班级、课程、学员、考试和训练记录。",
            "强调组织管理能力和全局视角。",
        ],
        "教官移动端价值",
        [
            "承担现场带班、签到签退、班级公告和课表查看等高频操作。",
            "适合培训现场快速进入班级、处理签到、接收提醒。",
            "强调执行效率和随时可达。",
        ],
    )
    capability_slide(
        prs,
        "教官网页端功能",
        "适合讲“管理与深度操作”",
        "6 / 14",
        "网页端能力",
        [
            "工作台：查看近期培训班、课程和考试待办，快速进入业务页面。",
            "培训班管理：查看班级列表、班级详情、学员名单、考试安排和变更记录。",
            "排课与周训练计划：查看培训班课表、训练计划和日程安排。",
            "教学资源：课程资源、资源库、我的资源、教学资源生成。",
            "考试协同：查看考试安排、成绩情况和班级考试结果。",
            "训练历史：查看班级历史记录，支撑后续复盘和个性化训练建议。",
        ],
        "汇报表达建议",
        [
            "把网页端讲成教官的“组织管理入口”，而不是简单的后台页面集合。",
            "突出其适合做完整配置、深度查看和全局统筹。",
            "强调它与移动端不是替代关系，而是前后协同关系。",
        ],
    )
    capability_slide(
        prs,
        "教官移动端功能",
        "适合讲“现场带班与即时执行”",
        "7 / 14",
        "移动端能力",
        [
            "首页工作台：查看个人信息、统计卡片、当日或本周安排。",
            "我的班级：快速进入班级列表和班级详情。",
            "日历课表：查看课程和考试安排，便于现场掌握节奏。",
            "当前课次管理：查看课程、教官、地点、状态等关键信息。",
            "签到签退：发起签到、结束签到、发起签退、查看实时结果。",
            "班级公告与通知：现场发布通知、查看未读提醒。",
        ],
        "汇报表达建议",
        [
            "把移动端讲成教官的“带班入口”和“现场执行入口”。",
            "重点突出签到签退、班级查看、公告提醒这些真实高频动作。",
            "这样能比单纯讲页面更容易体现产品落地价值。",
        ],
        accent=GREEN,
    )
    role_overview_slide(
        prs,
        "学员角色业务全景",
        "学员的核心目标：知道参加什么班、什么时候上课、如何签到、学什么、怎么考试",
        "8 / 14",
        "学员网页端价值",
        [
            "承担在线学习、考试参与、训练历史查看等连续性操作。",
            "适合完成课程学习、考试作答和个人记录查看。",
            "强调学习考试与结果沉淀。",
        ],
        "学员移动端价值",
        [
            "承担参训执行、扫码签到、班级查看和碎片化学习等高频操作。",
            "适合随时随地查看班级、课表、通知和资源。",
            "强调参训便捷性和现场执行效率。",
        ],
    )
    capability_slide(
        prs,
        "学员网页端功能",
        "适合讲“学习、考试、记录查看”",
        "9 / 14",
        "网页端能力",
        [
            "学员工作台：查看学习统计、考试次数和近期安排。",
            "培训班参与：查看班级、进入详情、发起报名申请。",
            "在线学习：查看课程资源、课程详情和学习进度。",
            "在线考试：查看考试列表、考试概览、在线作答和考试结果。",
            "训练历史：查看个人参训记录和考试记录。",
            "个人中心：沉淀个人信息和训练画像。",
        ],
        "汇报表达建议",
        [
            "把网页端讲成学员的“学习考试入口”，更适合长时学习和完整操作。",
            "重点体现从学习到考试再到结果查看的一条链路。",
            "这样能自然衔接到训练历史和个人档案价值。",
        ],
    )
    capability_slide(
        prs,
        "学员移动端功能",
        "适合讲“参训执行与碎片化学习”",
        "10 / 14",
        "移动端能力",
        [
            "首页：查看个人信息、统计卡片和快捷入口。",
            "我的班级：查看参加中的班级和班级详情。",
            "日历课表：按周查看训练与考试安排。",
            "签到签退：扫码签到、扫码签退或直接签到。",
            "移动学习：查看课程资源、资源库、推荐资源和资源详情。",
            "移动考试：查看考试、参加考试、查看考试结果和通知提醒。",
        ],
        "汇报表达建议",
        [
            "把移动端讲成学员的“参训入口”，核心是签到、查看安排和移动学习。",
            "强调它贴近真实训练现场，也能覆盖学习和考试的碎片化场景。",
            "这样更容易让听众理解移动端的必要性。",
        ],
        accent=ORANGE,
    )
    scene_slide(prs)
    support_slide(prs)
    result_slide(prs)
    next_step_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    output = Path(r"d:\QWQ\feishu\智慧教育训练平台介绍-角色终端版.pptx")
    build_deck(output)
    print(output)
