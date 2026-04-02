from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(r"d:\playform\police-training-platform\tmp_shuke_style_assets")
OUTPUT = Path(r"d:\QWQ\feishu\智慧教育训练平台介绍-角色终端版-授客风格.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(30, 39, 76)
DARK = RGBColor(22, 28, 54)
BLUE = RGBColor(104, 138, 255)
TEAL = RGBColor(110, 220, 212)
PURPLE = RGBColor(122, 132, 255)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(244, 246, 252)
TEXT = RGBColor(56, 64, 84)
MUTED = RGBColor(121, 130, 149)
LINE = RGBColor(226, 230, 240)


def rgb(t):
    return (t[0], t[1], t[2], 255)


def add_glow(base, xy, radius, color, blur):
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    x, y = xy
    draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=color)
    overlay = overlay.filter(ImageFilter.GaussianBlur(blur))
    return Image.alpha_composite(base, overlay)


def make_cover_bg(path: Path):
    w, h = 1600, 900
    base = Image.new("RGBA", (w, h), (34, 40, 76, 255))
    base = add_glow(base, (250, 210), 220, (72, 110, 220, 110), 110)
    base = add_glow(base, (1250, 190), 180, (96, 218, 208, 90), 100)
    base = add_glow(base, (1180, 640), 220, (91, 118, 255, 120), 110)
    base = add_glow(base, (840, 120), 90, (255, 126, 242, 180), 40)
    base = add_glow(base, (840, 120), 60, (255, 255, 255, 220), 15)

    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)

    # large dark glass panel hint
    d.rounded_rectangle((250, 110, 1210, 545), radius=32, fill=(255, 255, 255, 22), outline=(255, 255, 255, 30), width=2)
    d.rounded_rectangle((960, 320, 1220, 615), radius=28, fill=(255, 255, 255, 28), outline=(255, 255, 255, 28), width=2)

    # bottom wave
    d.ellipse((-220, 510, 1820, 1330), fill=(245, 247, 252, 255))
    d.ellipse((-260, 420, 1860, 1105), outline=(109, 218, 212, 255), width=22)
    d.ellipse((-250, 430, 1870, 1115), outline=(102, 136, 255, 200), width=12)
    base = Image.alpha_composite(base, overlay)
    base.convert("RGB").save(path)


def make_dark_section_bg(path: Path):
    w, h = 1600, 900
    base = Image.new("RGBA", (w, h), (25, 29, 50, 255))
    base = add_glow(base, (180, 180), 240, (43, 115, 198, 95), 120)
    base = add_glow(base, (1360, 250), 220, (106, 60, 136, 100), 115)
    base = add_glow(base, (1280, 700), 200, (76, 124, 255, 90), 110)
    base = add_glow(base, (350, 680), 170, (214, 86, 98, 48), 90)

    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    for y in range(120, 760, 80):
        d.line((860, y, 1160, y - 65), fill=(255, 132, 132, 55), width=2)
        d.line((700, y, 1020, y + 80), fill=(83, 124, 255, 40), width=2)
    for x in range(220, 1500, 220):
        d.ellipse((x, 120 + (x % 170), x + 8, 128 + (x % 170)), fill=(255, 90, 90, 90))
    base = Image.alpha_composite(base, overlay)
    base.convert("RGB").save(path)


def make_light_bg(path: Path):
    w, h = 1600, 900
    base = Image.new("RGBA", (w, h), (245, 247, 252, 255))
    base = add_glow(base, (90, 820), 280, (122, 229, 220, 85), 120)
    base = add_glow(base, (1510, 100), 260, (123, 185, 255, 65), 120)
    base = add_glow(base, (1380, 650), 240, (123, 134, 255, 95), 95)

    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    d.ellipse((1225, 380, 1790, 980), fill=(255, 255, 255, 185))
    d.ellipse((1260, 420, 1750, 940), outline=(139, 151, 255, 35), width=20)
    d.ellipse((1290, 450, 1710, 910), outline=(139, 151, 255, 30), width=16)
    base = Image.alpha_composite(base, overlay)
    base.convert("RGB").save(path)


def build_assets():
    OUT_DIR.mkdir(exist_ok=True)
    make_cover_bg(OUT_DIR / "cover.png")
    make_dark_section_bg(OUT_DIR / "section.png")
    make_light_bg(OUT_DIR / "light.png")


def add_picture_bg(slide, img_path):
    slide.shapes.add_picture(str(img_path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_shadow_card(slide, left, top, width, height, transparency=0.0):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.04), top + Inches(0.06), width, height, RGBColor(204, 212, 230), transparency=0.55)
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, WHITE, line=LINE, transparency=transparency)


def add_textbox(slide, left, top, width, height, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, italic=False, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=16, color=TEXT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.left_margin = Pt(18)
        p.hanging = Pt(-10)
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return tb


def add_logo(slide):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.62), Inches(0.26), Inches(0.38), Inches(0.38), RGBColor(61, 177, 114))
    add_textbox(slide, Inches(11.67), Inches(0.335), Inches(0.28), Inches(0.15), "训", 14, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(12.05), Inches(0.28), Inches(0.90), Inches(0.18), "智慧训练", 11, RGBColor(101, 108, 128), True)


def add_center_title(slide, title, subtitle=None):
    add_logo(slide)
    add_textbox(slide, Inches(2.4), Inches(0.82), Inches(8.5), Inches(0.35), title, 20, RGBColor(103, 200, 214), True, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.05), Inches(1.28), Inches(0.45), Inches(0.02), BLUE)
    if subtitle:
        add_textbox(slide, Inches(1.2), Inches(1.45), Inches(10.9), Inches(0.45), subtitle, 13.5, MUTED, False, PP_ALIGN.CENTER)


def add_pill(slide, left, top, width, text, fill):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.34), fill)
    add_textbox(slide, left + Inches(0.05), top + Inches(0.06), width - Inches(0.10), Inches(0.16), text, 10.5, WHITE, True, PP_ALIGN.CENTER)


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "cover.png")
    add_textbox(slide, Inches(2.12), Inches(2.65), Inches(9.1), Inches(0.8), "智慧教育训练平台", 33, RGBColor(117, 217, 224), True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.75), Inches(3.35), Inches(7.8), Inches(0.45), "教官 / 学员角色汇报", 21, RGBColor(117, 140, 255), True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.7), Inches(4.05), Inches(3.8), Inches(0.22), "网页端 + 移动端协同视角", 13, RGBColor(221, 226, 239), False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.55), Inches(6.55), Inches(2.3), Inches(0.20), "项目汇报材料", 11, RGBColor(86, 98, 128), False, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.20), Inches(4.7), Inches(0.08), Inches(0.08), WHITE)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.38), Inches(4.7), Inches(0.08), Inches(0.08), RGBColor(217, 221, 231))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.56), Inches(4.7), Inches(0.08), Inches(0.08), RGBColor(217, 221, 231))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.74), Inches(4.7), Inches(0.08), Inches(0.08), RGBColor(217, 221, 231))
    return slide


def contents_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "cover.png")
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(4.02), SLIDE_W, Inches(3.48), RGBColor(245, 247, 252))
    add_textbox(slide, Inches(4.2), Inches(1.3), Inches(4.9), Inches(0.6), "CONTENTS", 34, RGBColor(111, 154, 255), False, PP_ALIGN.CENTER, italic=True, font="Calibri")

    cards = [
        ("01", "双端定位", ["角色 × 终端", "汇报逻辑"]),
        ("02", "教官端能力", ["网页端管理", "移动端带班"]),
        ("03", "学员端能力", ["网页端学习考试", "移动端参训"]),
        ("04", "双端协同", ["典型场景", "统一支撑"]),
        ("05", "成果与规划", ["当前成果", "下一步建议"]),
    ]
    start = Inches(0.9)
    for idx, (num, title, bullets) in enumerate(cards):
        x = start + Inches(idx * 2.47)
        add_shadow_card(slide, x, Inches(3.32), Inches(1.98), Inches(2.05), 0.15)
        add_textbox(slide, x + Inches(0.52), Inches(3.55), Inches(0.95), Inches(0.32), num, 22, WHITE, False, PP_ALIGN.CENTER, italic=True, font="Calibri")
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.32), Inches(1.98), Inches(0.54), RGBColor(255, 255, 255), transparency=0.35)
        add_textbox(slide, x + Inches(0.12), Inches(4.15), Inches(1.74), Inches(0.25), title, 15.5, TEXT, True, PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.18), Inches(4.62), Inches(1.64), Inches(0.52), "\n".join([f"· {b}" for b in bullets]), 11.5, MUTED, False, PP_ALIGN.LEFT)
    return slide


def light_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "light.png")
    add_center_title(slide, "同一套训练业务的双端定位", "网页端负责管理与深度操作，移动端负责现场执行与高频使用")
    cards = [
        ("教官网页端", "培训班组织\n课程与资源\n考试协同", TEAL),
        ("教官移动端", "班级查看\n签到签退\n现场通知", BLUE),
        ("学员网页端", "在线学习\n考试参与\n训练记录", TEAL),
        ("学员移动端", "我的班级\n扫码签到\n移动学习", BLUE),
    ]
    x_positions = [Inches(0.95), Inches(3.95), Inches(6.95), Inches(9.95)]
    for (title, body, color), x in zip(cards, x_positions):
        add_shadow_card(slide, x, Inches(2.45), Inches(2.35), Inches(3.15), 0.0)
        add_pill(slide, x + Inches(0.35), Inches(2.18), Inches(1.65), title, color)
        add_textbox(slide, x + Inches(0.20), Inches(2.95), Inches(1.95), Inches(1.6), body, 15, TEXT, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.0), Inches(6.28), Inches(9.2), Inches(0.30), "后续汇报不再把两端拆成两套系统，而是围绕同一条训练业务链路，看不同角色在不同终端上的使用方式。", 13.5, MUTED, False, PP_ALIGN.CENTER)
    return slide


def loop_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "light.png")
    add_center_title(slide, "角色与终端协同下的核心业务闭环", "围绕训练组织、学习、签到、考试和结果沉淀形成一体化闭环")
    steps = [
        "人员进入系统",
        "培训班组织",
        "课程与资源",
        "现场签到",
        "在线考试",
        "训练记录",
        "结果分析",
    ]
    x = Inches(0.7)
    widths = [1.55, 1.65, 1.6, 1.45, 1.45, 1.55, 1.45]
    colors = [TEAL, BLUE, PURPLE, TEAL, BLUE, PURPLE, TEAL]
    for idx, (step, w, color) in enumerate(zip(steps, widths, colors)):
        add_shadow_card(slide, x, Inches(3.18), Inches(w), Inches(1.22), 0.0)
        add_pill(slide, x + Inches(0.20), Inches(3.00), Inches(0.65), f"0{idx+1}", color)
        add_textbox(slide, x + Inches(0.10), Inches(3.62), Inches(w) - Inches(0.20), Inches(0.34), step, 13.5, TEXT, True, PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_textbox(slide, x + Inches(w) + Inches(0.03), Inches(3.54), Inches(0.26), Inches(0.24), "→", 24, BLUE, True, PP_ALIGN.CENTER, font="Calibri")
        x += Inches(w) + Inches(0.17)
    add_shadow_card(slide, Inches(1.15), Inches(5.35), Inches(4.85), Inches(0.92), 0.0)
    add_pill(slide, Inches(1.42), Inches(5.16), Inches(1.55), "网页端承担", BLUE)
    add_textbox(slide, Inches(1.40), Inches(5.66), Inches(4.3), Inches(0.26), "组织、配置、查看、分析", 16, TEXT, True, PP_ALIGN.CENTER)
    add_shadow_card(slide, Inches(7.25), Inches(5.35), Inches(4.85), Inches(0.92), 0.0)
    add_pill(slide, Inches(7.52), Inches(5.16), Inches(1.55), "移动端承担", TEAL)
    add_textbox(slide, Inches(7.50), Inches(5.66), Inches(4.3), Inches(0.26), "签到、查看、执行、提醒", 16, TEXT, True, PP_ALIGN.CENTER)
    return slide


def dark_section_slide(prs, english, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "section.png")
    add_textbox(slide, Inches(4.4), Inches(1.85), Inches(4.6), Inches(0.5), english, 24, RGBColor(112, 149, 255), False, PP_ALIGN.CENTER, italic=True, font="Calibri")
    add_textbox(slide, Inches(2.35), Inches(3.0), Inches(8.7), Inches(0.6), title, 28, WHITE, True, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.08), Inches(3.78), Inches(0.46), Inches(0.02), RGBColor(214, 222, 246))
    add_textbox(slide, Inches(2.0), Inches(4.18), Inches(9.4), Inches(0.6), subtitle, 16, RGBColor(226, 232, 244), False, PP_ALIGN.CENTER)
    return slide


def dual_cards_slide(prs, title, subtitle, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "light.png")
    add_center_title(slide, title, subtitle)
    add_shadow_card(slide, Inches(0.85), Inches(2.15), Inches(5.7), Inches(4.15), 0.0)
    add_pill(slide, Inches(1.18), Inches(1.96), Inches(1.85), left_title, TEAL)
    add_bullets(slide, Inches(1.1), Inches(2.68), Inches(5.0), Inches(3.2), left_items, 15.5)
    add_shadow_card(slide, Inches(6.78), Inches(2.15), Inches(5.7), Inches(4.15), 0.0)
    add_pill(slide, Inches(7.10), Inches(1.96), Inches(1.85), right_title, BLUE)
    add_bullets(slide, Inches(7.02), Inches(2.68), Inches(5.0), Inches(3.2), right_items, 15.5)
    return slide


def four_card_slide(prs, title, subtitle, cards):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "light.png")
    add_center_title(slide, title, subtitle)
    positions = [
        (Inches(0.92), Inches(2.2)),
        (Inches(3.97), Inches(2.2)),
        (Inches(7.02), Inches(2.2)),
        (Inches(10.07), Inches(2.2)),
    ]
    fills = [TEAL, BLUE, PURPLE, BLUE]
    for (header, body), (left, top), fill in zip(cards, positions, fills):
        add_shadow_card(slide, left, top, Inches(2.35), Inches(3.55), 0.0)
        add_pill(slide, left + Inches(0.34), top - Inches(0.18), Inches(1.55), header, fill)
        add_textbox(slide, left + Inches(0.18), top + Inches(0.72), Inches(1.98), Inches(2.2), body, 14.5, TEXT, False, PP_ALIGN.CENTER)
    return slide


def scene_list_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "light.png")
    add_center_title(slide, "双端协同典型场景", "不是两套系统各讲各的，而是同一条训练链路在不同终端间自然衔接")
    items = [
        ("培训班组织", "教官网页端完成班级组织、课表安排和资源准备，学员通过网页端或移动端查看班级与日历。", TEAL),
        ("现场签到", "教官移动端发起签到签退，学员移动端扫码完成签到，结果实时回流到班级记录。", BLUE),
        ("学习支撑", "教官网页端准备课程资源与教学内容，学员网页端和移动端持续查看课程与资源。", PURPLE),
        ("考试闭环", "教官网页端协同考试安排，学员网页端或移动端参加考试，成绩统一沉淀到训练记录。", BLUE),
    ]
    top = Inches(2.05)
    for idx, (head, desc, fill) in enumerate(items):
        y = top + Inches(idx * 1.1)
        add_shadow_card(slide, Inches(1.1), y, Inches(10.9), Inches(0.8), 0.0)
        add_pill(slide, Inches(1.32), y + Inches(0.18), Inches(1.25), head, fill)
        add_textbox(slide, Inches(2.85), y + Inches(0.22), Inches(8.65), Inches(0.32), desc, 13.2, MUTED)
    return slide


def support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "light.png")
    add_center_title(slide, "统一支撑关系", "统一角色、统一权限、统一训练数据，支撑网页端和移动端的协同运行")
    add_shadow_card(slide, Inches(1.05), Inches(1.95), Inches(11.2), Inches(0.95), 0.0)
    add_textbox(slide, Inches(1.35), Inches(2.26), Inches(10.6), Inches(0.26), "同一套后端与数据体系，支撑教官和学员在网页端、移动端的不同使用方式。", 17, TEXT, True, PP_ALIGN.CENTER)
    labels = [
        ("教官网页端", TEAL),
        ("教官移动端", BLUE),
        ("学员网页端", PURPLE),
        ("学员移动端", BLUE),
    ]
    xs = [Inches(1.2), Inches(4.1), Inches(7.0), Inches(9.9)]
    for (text, fill), x in zip(labels, xs):
        add_shadow_card(slide, x, Inches(4.5), Inches(2.25), Inches(1.22), 0.0)
        add_pill(slide, x + Inches(0.35), Inches(4.30), Inches(1.55), text, fill)
    center_modules = [("训练组织", TEAL), ("课程资源", BLUE), ("在线考试", PURPLE), ("训练记录", TEAL)]
    module_x = Inches(1.6)
    for text, fill in center_modules:
        add_shadow_card(slide, module_x, Inches(3.15), Inches(2.3), Inches(0.9), 0.0)
        add_textbox(slide, module_x + Inches(0.15), Inches(3.48), Inches(2.0), Inches(0.18), text, 15, TEXT, True, PP_ALIGN.CENTER)
        module_x += Inches(2.55)
    return slide


def result_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "light.png")
    add_center_title(slide, "当前成果与演示重点", "聚焦当前项目中已经具备展示价值的能力，突出双端落地场景")
    add_shadow_card(slide, Inches(0.9), Inches(2.1), Inches(5.75), Inches(4.2), 0.0)
    add_pill(slide, Inches(1.2), Inches(1.92), Inches(1.7), "网页端成果", BLUE)
    add_bullets(
        slide,
        Inches(1.12),
        Inches(2.65),
        Inches(5.0),
        Inches(3.2),
        [
            "教官网页端已覆盖工作台、培训班管理、排课查看、课程资源、资源库、考试协同等核心能力。",
            "学员网页端已覆盖学习、考试、报名、训练历史和个人中心等关键链路。",
            "适合突出“管理与深度操作”的桌面端价值。",
        ],
        15.5,
    )
    add_shadow_card(slide, Inches(6.8), Inches(2.1), Inches(5.65), Inches(4.2), 0.0)
    add_pill(slide, Inches(7.1), Inches(1.92), Inches(1.7), "移动端成果", TEAL)
    add_bullets(
        slide,
        Inches(7.02),
        Inches(2.65),
        Inches(4.95),
        Inches(3.2),
        [
            "移动端已覆盖首页工作台、我的班级、日历课表、签到签退、资源查看、在线考试和通知入口。",
            "适合突出“带班现场、扫码签到、移动参训”的高频执行价值。",
            "更容易体现项目在真实训练场景中的落地性。",
        ],
        15.5,
    )
    return slide


def next_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_bg(slide, OUT_DIR / "light.png")
    add_center_title(slide, "下一步优化建议", "在当前角色 × 终端结构上，继续增强训练闭环的完整性和汇报表达的清晰度")
    items = [
        ("继续补强角色页细节", "完善报名、通知、个人中心等角色端细节体验。", TEAL),
        ("强化双端消息联动", "围绕公告、考试提醒、签到异常和待办做统一触达。", BLUE),
        ("提升结果沉淀能力", "把签到、学习、考试进一步沉淀到个人训历与评估分析。", PURPLE),
        ("统一汇报话术", "始终围绕角色 × 终端 × 协同闭环展开，不再回到页面堆砌式表达。", BLUE),
    ]
    positions = [
        (Inches(1.0), Inches(2.15)),
        (Inches(6.9), Inches(2.15)),
        (Inches(1.0), Inches(4.35)),
        (Inches(6.9), Inches(4.35)),
    ]
    for (head, desc, fill), (left, top) in zip(items, positions):
        add_shadow_card(slide, left, top, Inches(5.25), Inches(1.55), 0.0)
        add_pill(slide, left + Inches(0.25), top + Inches(0.18), Inches(1.95), head, fill)
        add_textbox(slide, left + Inches(0.26), top + Inches(0.78), Inches(4.6), Inches(0.44), desc, 13.2, MUTED)
    return slide


def build_ppt():
    build_assets()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_slide(prs)
    contents_slide(prs)
    light_intro_slide(prs)
    loop_slide(prs)
    dark_section_slide(prs, "INSTRUCTOR", "教官角色业务全景", "网页端负责组织与管理，移动端负责带班与现场执行。")
    dual_cards_slide(
        prs,
        "教官端能力说明",
        "网页端强调管理深度，移动端强调带班效率",
        "教官网页端",
        [
            "工作台：查看近期培训班、课程和考试待办。",
            "培训班管理：查看班级详情、学员名单、考试安排和变更记录。",
            "排课与周训练计划：查看培训班课表和训练计划。",
            "教学资源：课程资源、资源库、我的资源、教学资源生成。",
            "考试协同：查看考试安排和班级考试结果。",
        ],
        "教官移动端",
        [
            "首页工作台：查看个人信息、统计卡片和近期安排。",
            "我的班级：快速进入班级列表和班级详情。",
            "日历课表：随时查看课程和考试安排。",
            "签到签退：发起签到、结束签到、发起签退。",
            "班级公告与通知：现场查看并处理通知提醒。",
        ],
    )
    four_card_slide(
        prs,
        "教官端汇报重点",
        "把教官端讲成“组织管理 + 现场带班”的一体化入口",
        [
            ("组织管理", "适合讲培训班组织、课程安排、资源准备和考试协同。"),
            ("深度查看", "适合讲班级详情、学员名单、变更记录和训练计划。"),
            ("现场带班", "适合讲签到签退、班级公告、课表查看和现场执行。"),
            ("双端协同", "强调网页端负责准备，移动端负责执行，形成完整闭环。"),
        ],
    )
    dark_section_slide(prs, "STUDENT", "学员角色业务全景", "网页端负责学习与考试，移动端负责参训与现场执行。")
    dual_cards_slide(
        prs,
        "学员端能力说明",
        "网页端强调学习考试，移动端强调参训签到",
        "学员网页端",
        [
            "学员工作台：查看学习统计、考试次数和近期安排。",
            "培训班参与：查看班级、进入详情、发起报名申请。",
            "在线学习：查看课程资源、课程详情和学习进度。",
            "在线考试：查看考试列表、考试概览、在线作答和考试结果。",
            "训练历史：查看个人参训记录和考试记录。",
        ],
        "学员移动端",
        [
            "首页：查看个人信息、统计卡片和快捷入口。",
            "我的班级：查看参加中的班级和班级详情。",
            "日历课表：按周查看训练与考试安排。",
            "签到签退：扫码签到、扫码签退或直接签到。",
            "移动学习与移动考试：随时查看资源、参加考试和接收提醒。",
        ],
    )
    four_card_slide(
        prs,
        "学员端汇报重点",
        "把学员端讲成“学习入口 + 参训入口”的协同体系",
        [
            ("学习入口", "网页端更适合做连续性学习和完整考试操作。"),
            ("参训入口", "移动端更适合查看班级、扫码签到和现场执行。"),
            ("记录沉淀", "学习、考试和签到最终都会沉淀到个人训练记录中。"),
            ("体验价值", "让学员在不同场景下都能顺畅完成参训和学习。"),
        ],
    )
    scene_list_slide(prs)
    support_slide(prs)
    result_slide(prs)
    next_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build_ppt()
